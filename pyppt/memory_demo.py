
from io import BytesIO, StringIO
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

prs_template = Presentation("./target/demo_tp.pptx")
target_stream = BytesIO()
prs_template.save(target_stream)

prs_output = Presentation(target_stream)

def replace_text_preserving_style(text_frame, new_text):
    """
    替换文本框内容，同时保留原有的文本样式。

    :param text_frame: pptx.text.text.TextFrame, 要替换内容的文本框。
    :param new_text: str, 新的文本内容。
    """
    paragraphs = text_frame.paragraphs
    current_index = 0

    for paragraph in paragraphs:
        for run in paragraph.runs:
            if current_index < len(new_text):
                # 按字符替换，保留样式
                run.text = new_text[current_index:current_index + len(run.text)]
                current_index += len(run.text)
            else:
                # 新文本用完，清空剩余内容
                run.text = ""
        
        # 如果新文本还有剩余，添加额外的 Run
        if current_index < len(new_text):
            remaining_text = new_text[current_index:]
            paragraph.add_run().text = remaining_text
            current_index += len(remaining_text)

def update_slide_content_preserving_style(slide, updates):
    """
    修改幻灯片中形状的内容（文本或图片），保留原样式。

    :param slide: pptx.slide.Slide, 要修改的幻灯片对象
    :param updates: dict, 要修改的内容，形状名称到新内容的映射，如：
        {
            "TitleBox": "New Title Text",
            "ImagePlaceholder": "path/to/new_image.jpg"
        }
    """
    for shape in slide.shapes:
        if shape.name in updates:
            new_content = updates[shape.name]
            
            # 如果是文本框，逐段修改文本内容，保留样式
            if shape.has_text_frame:
                replace_text_preserving_style(shape.text_frame, new_content) # 替换内容，仅修改文本不影响样式
            
            # 如果是图片占位符，替换图片，保留位置和尺寸
            elif shape.shape_type == 13:  # 13 表示图片形状
                with open(new_content, "rb") as image_file:
                    # 保留原始位置和尺寸
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    # 删除旧图片
                    shape.element.getparent().remove(shape.element)
                    # 插入新图片
                    slide.shapes.add_picture(image_file, left, top, width, height)


import copy

def duplicate_shape(target_slide, source_shape):
    """
    复制单个形状到目标幻灯片。
    
    :param target_slide: 目标幻灯片对象
    :param source_shape: 源形状对象
    """
    if source_shape.shape_type == 13:  # 图片形状
        image_stream = BytesIO(source_shape.image.blob)
        target_slide.shapes.add_picture(
            image_stream, source_shape.left, source_shape.top, source_shape.width, source_shape.height
        )
    elif source_shape.shape_type == 6:  # 组合形状
        # 创建新的组合形状
        group = target_slide.shapes.add_group_shape()
        group.left, group.top, group.width, group.height = source_shape.left, source_shape.top, source_shape.width, source_shape.height
        
        # 递归复制子形状
        for sub_shape in source_shape.shapes:
            duplicate_shape(group, sub_shape)
    else:
        # 对于其他类型的形状，直接复制 XML
        el = source_shape.element
        new_el = copy.deepcopy(el)
        target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')


def copy_notes(source_slide, target_slide):
    """
    复制备注内容，从源幻灯片到目标幻灯片。
    """
    if source_slide.has_notes_slide:
        # 确保目标幻灯片也有 notes_slide
        if not target_slide.has_notes_slide:
            target_slide.notes_slide  # 初始化 notes_slide
        
        source_notes = source_slide.notes_slide
        target_notes = target_slide.notes_slide

        # 遍历源备注中的形状
        for shape, target_shape in zip(source_notes.shapes, target_notes.shapes):
            if shape.has_text_frame and target_shape.has_text_frame:
                target_shape.text_frame.text = shape.text_frame.text



def duplicate_slide(pres, index, copy_name = None):
    """
    复制 PowerPoint 幻灯片，包括图片、形状和备注。
    
    :param pres: Presentation 对象
    :param index: 要复制的幻灯片索引
    :return: 新幻灯片对象
    """
    source_slide = pres.slides[index]
    blank_slide_layout = pres.slide_layouts[6]  # 通常使用空白布局
    new_slide = pres.slides.add_slide(blank_slide_layout)

        # 复制所有形状
    for shape in source_slide.shapes:
        if copy_name and copy_name != shape.name:
            continue
        duplicate_shape(new_slide, shape)

    copy_notes(source_slide, new_slide)
    
    return new_slide


# 要修改的内容
updates = {
    "s_title": "Updated Title Text",
    "s_content":     """
        恶魔 神佛
    """,
    # "p_s1": "./pyppt/good.png"
}

slide = prs_output.slides[0]
# 修改幻灯片内容
update_slide_content_preserving_style(slide, updates)

duplicate_slide(prs_output, 0)
duplicate_slide(prs_output, 0, "g_2")
file_path = "./target/demo_output.pptx"
prs_output.save(file_path)


import os

# 检查文件是否存在
if os.path.exists(file_path):
    # 根据操作系统打开文件
    if os.name == 'nt':  # Windows
        os.startfile(file_path)
    elif os.name == 'posix':  # macOS 或 Linux
        os.system(f'open "{file_path}"')  # macOS 使用 'open'
        # os.system(f'xdg-open "{file_path}"')  # Linux 通常使用 'xdg-open'
else:
    print("文件不存在，请检查路径！")