import random
from pptx import Presentation
from pptx.util import Pt, Inches

from PIL import Image

import os

# 当前 Python 文件的绝对路径
current_file_path = os.path.abspath(__file__)

# 当前文件所在的目录
current_dir = os.path.dirname(current_file_path)


def _r(fp):
    return os.path.join(current_dir, fp)

# 打开 WEBP 图片
image = Image.open(_r("./good.webp"))

# 保存为 PNG 格式
image.save(_r("./good.png"), "PNG")



# 名片信息
cards = [
    {"name": "John Doe", "title": "Software Engineer", "phone": "123-456-7890", "email": "john.doe@example.com"},
    {"name": "Jane Smith", "title": "Data Scientist", "phone": "987-654-3210", "email": "jane.smith@example.com"},
    {"name": "Alice Johnson", "title": "Product Manager", "phone": "555-666-7777", "email": "alice.johnson@example.com"},
    {"name": "Alice Johnson22", "title": "Product Manager", "phone": "555-666-7777", "email": "alice.johnson@example.com"},
]

# 头像文件路径
avatars = [
    _r("./good.png"),
    _r("./good.png"),
    _r("./good.png"),
    _r("./good.png"),
]

from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE


# 创建PPT
prs = Presentation()

for i, card in enumerate(cards):
    slide_layout = prs.slide_layouts[1]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)

    chart_data = ChartData()
    chart_data.categories = ['Q1 Sales', 'Q2 Sales', 'Q3 Sales']
    chart_data.add_series('West',    (32.2, 28.4, 34.7))
    chart_data.add_series('East',    (24.3, 30.6, 20.2))
    chart_data.add_series('Midwest', (20.4, 18.3, 26.2))

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.series[0].smooth = True

    # 添加头像
    slide.shapes.add_picture(avatars[i], Inches(0.5), Inches(0.5), Inches(2), Inches(2))
    slide.shapes.add_picture(avatars[i], Inches(0.5), Inches(4.5), Inches(2), Inches(2))


    # 添加姓名
    title_box = slide.shapes.add_textbox(Inches(3), Inches(0.5), Inches(6), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = card["name"]
    title_frame.paragraphs[0].font.size = Pt(32)

    # 添加职位
    title_box = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(6), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = f"Title: {card['title']}"
    title_frame.paragraphs[0].font.size = Pt(20)

    # 添加联系方式
    contact_box = slide.shapes.add_textbox(Inches(3), Inches(2.5), Inches(6), Inches(2))
    contact_frame = contact_box.text_frame
    contact_frame.text = f"Phone: {card['phone']}\nEmail: {card['email']}"
    for paragraph in contact_frame.paragraphs:
        paragraph.font.size = Pt(18)

# 保存PPT
prs.save(_r("business_cards_with_avatars.pptx"))

# 成语列表
idioms = [
    "一心一意", "百发百中", "千言万语", "万事如意", 
    "十全十美", "百折不挠", "画龙点睛", "锦上添花"
]


# 读取PPT并修改内容
def add_idioms_to_ppt(filename):
    prs = Presentation(filename)

    for i, slide in enumerate(prs.slides):
        # 随机选择一个成语
        idiom = random.choice(idioms)

        slide.shapes.title.text = f"\n\nGOOD: {idiom}"

        # 在每页内容后添加成语
        # for shape in slide.shapes:
        #     if shape.has_text_frame:
        #         shape.text += f"\n\n成语: {idiom}"

    prs.save(filename)

add_idioms_to_ppt(_r("business_cards_with_avatars.pptx"))